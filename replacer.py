# -*- coding: utf-8 -*-
"""
文档词汇批量替换 - 核心逻辑模块
"""

import os
import csv
import shutil
import re
from pathlib import Path
from typing import List, Tuple

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

# 支持的文件扩展名
SUPPORTED_EXTENSIONS = [".docx", ".txt", ".xlsx", ".xls", ".pptx", ".ppt", ".csv", ".doc"]


# ─── 替换规则加载 ────────────────────────────────────────────

def load_rules(filepath: str) -> List[Tuple[str, str]]:
    """
    从 xlsx 或 csv 文件加载替换规则。
    返回 [(原词, 替换词), ...] 列表。
    """
    ext = Path(filepath).suffix.lower()
    if ext == ".xlsx":
        return _load_rules_xlsx(filepath)
    elif ext == ".csv":
        return _load_rules_csv(filepath)
    else:
        raise ValueError(f"不支持的文件格式: {ext}，请使用 .xlsx 或 .csv 文件")


def _load_rules_xlsx(filepath: str) -> List[Tuple[str, str]]:
    wb = load_workbook(filepath, read_only=True)
    ws = wb.active
    rules = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i == 0:
            # 跳过标题行（如果第一行看起来像标题）
            if row and isinstance(row[0], str) and row[0].strip() in ("原词", "原文", "查找", "原始"):
                continue
        if row and len(row) >= 2 and row[0] and row[1]:
            rules.append((str(row[0]).strip(), str(row[1]).strip()))
    wb.close()
    return rules


def _load_rules_csv(filepath: str) -> List[Tuple[str, str]]:
    rules = []
    # 尝试检测编码
    for encoding in ("utf-8-sig", "utf-8", "gbk", "gb2312", "latin1"):
        try:
            with open(filepath, "r", encoding=encoding) as f:
                reader = csv.reader(f)
                first = True
                for row in reader:
                    if first:
                        first = False
                        if row and row[0].strip() in ("原词", "原文", "查找", "原始"):
                            continue
                    if len(row) >= 2 and row[0].strip() and row[1].strip():
                        rules.append((row[0].strip(), row[1].strip()))
            return rules
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError("无法读取 CSV 文件，请确认文件编码为 UTF-8 或 GBK")


# ─── 备份与还原 ─────────────────────────────────────────────

BACKUP_DIR_NAME = "_backup"


def _get_backup_path(filepath: str) -> str:
    """获取备份文件的路径"""
    p = Path(filepath)
    backup_dir = p.parent / BACKUP_DIR_NAME
    return str(backup_dir / p.name)


def backup_file(filepath: str) -> str:
    """
    备份原文件到 _backup 文件夹。
    返回备份文件路径。
    """
    p = Path(filepath)
    backup_dir = p.parent / BACKUP_DIR_NAME
    backup_dir.mkdir(exist_ok=True)
    backup_path = backup_dir / p.name
    shutil.copy2(filepath, backup_path)
    return str(backup_path)


def restore_file(filepath: str) -> bool:
    """
    从备份还原文件。
    返回是否成功还原。
    """
    backup_path = _get_backup_path(filepath)
    if not os.path.exists(backup_path):
        return False
    shutil.copy2(backup_path, filepath)
    return True


def has_backup(filepath: str) -> bool:
    """检查文件是否有备份"""
    return os.path.exists(_get_backup_path(filepath))


# ─── 文档替换核心 ─────────────────────────────────────────────

def replace_in_docx(filepath: str, rules: List[Tuple[str, str]]) -> dict:
    """
    对 docx 文件执行批量替换。
    保留原始格式，处理跨 run 的文本匹配。

    返回 {"total_replacements": int, "detail": {原词: 替换次数}}
    """
    doc = Document(filepath)
    detail = {}

    for old_text, new_text in rules:
        count = 0
        # 替换段落中的文本
        for para in doc.paragraphs:
            count += _replace_in_paragraph(para, old_text, new_text)

        # 替换表格中的文本
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        count += _replace_in_paragraph(para, old_text, new_text)

        # 替换页眉页脚中的文本
        for section in doc.sections:
            for header_footer in [section.header, section.footer]:
                if header_footer is not None:
                    for para in header_footer.paragraphs:
                        count += _replace_in_paragraph(para, old_text, new_text)
                    for table in header_footer.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                for para in cell.paragraphs:
                                    count += _replace_in_paragraph(para, old_text, new_text)

        if count > 0:
            detail[old_text] = count

    doc.save(filepath)

    total = sum(detail.values())
    return {"total_replacements": total, "detail": detail}


def _replace_in_paragraph(paragraph, old_text: str, new_text: str) -> int:
    """
    在段落中替换文本，尽量保留原始格式。
    处理跨 run 的文本匹配。
    返回替换次数。
    """
    # 先尝试简单的单 run 替换
    count = 0
    for run in paragraph.runs:
        if old_text in run.text:
            run.text = run.text.replace(old_text, new_text)
            count += run.text.count(new_text)  # 近似计数
            # 重新计算：用替换前后的差异来准确计数
    if count > 0:
        return count

    # 如果单 run 没找到，尝试跨 run 匹配
    full_text = "".join(run.text for run in paragraph.runs)
    if old_text not in full_text:
        return 0

    # 跨 run 替换策略
    count = full_text.count(old_text)
    if count == 0:
        return 0

    return _cross_run_replace(paragraph, old_text, new_text)


def _cross_run_replace(paragraph, old_text: str, new_text: str) -> int:
    """
    处理跨 run 的文本替换。
    核心策略：构建字符到 run 的映射，找到匹配位置后，
    在第一个涉及的 run 中放入替换文本，清空其余涉及的 run 中对应的字符。
    """
    runs = paragraph.runs
    if not runs:
        return 0

    # 构建字符位置到 (run_index, char_index_in_run) 的映射
    char_map = []
    for run_idx, run in enumerate(runs):
        for char_idx in range(len(run.text)):
            char_map.append((run_idx, char_idx))

    full_text = "".join(run.text for run in runs)
    count = 0
    search_start = 0

    while True:
        pos = full_text.find(old_text, search_start)
        if pos == -1:
            break

        count += 1
        match_end = pos + len(old_text)

        # 找到匹配涉及的 run 范围
        first_run_idx = char_map[pos][0]
        last_run_idx = char_map[match_end - 1][0]

        # 在每个涉及的 run 中处理文本
        for run_idx in range(first_run_idx, last_run_idx + 1):
            run = runs[run_idx]
            run_text = run.text

            # 计算当前 run 中需要替换的字符范围
            # run 在全文中的起始位置
            run_start_in_full = sum(len(runs[i].text) for i in range(run_idx))
            run_end_in_full = run_start_in_full + len(run_text)

            # 匹配区域在当前 run 中的范围
            local_start = max(0, pos - run_start_in_full)
            local_end = min(len(run_text), match_end - run_start_in_full)

            if run_idx == first_run_idx:
                # 第一个 run：用替换文本替换匹配部分
                run.text = run_text[:local_start] + new_text + run_text[local_end:]
            else:
                # 后续 run：仅移除匹配部分
                run.text = run_text[:local_start] + run_text[local_end:]

        # 由于文本已变化，需要重建映射后继续（简单起见，递归处理剩余）
        # 这里简单 break 后再调用一次
        return count + _cross_run_replace(paragraph, old_text, new_text)

    return count


# ─── TXT 替换 ─────────────────────────────────────────────

def replace_in_txt(filepath: str, rules: List[Tuple[str, str]]) -> dict:
    """
    对 txt 文件执行批量替换。
    返回 {"total_replacements": int, "detail": {原词: 替换次数}}
    """
    # 尝试多种编码读取
    content = None
    detected_encoding = None
    for encoding in ("utf-8-sig", "utf-8", "gbk", "gb2312", "latin1"):
        try:
            with open(filepath, "r", encoding=encoding) as f:
                content = f.read()
            detected_encoding = encoding
            break
        except (UnicodeDecodeError, UnicodeError):
            continue

    if content is None:
        raise ValueError("无法读取文本文件，请确认文件编码")

    detail = {}
    for old_text, new_text in rules:
        count = content.count(old_text)
        if count > 0:
            content = content.replace(old_text, new_text)
            detail[old_text] = count

    with open(filepath, "w", encoding=detected_encoding) as f:
        f.write(content)

    total = sum(detail.values())
    return {"total_replacements": total, "detail": detail}


# ─── XLSX 替换 ────────────────────────────────────────────

def replace_in_xlsx(filepath: str, rules: List[Tuple[str, str]]) -> dict:
    """
    对 xlsx 文件执行批量替换。
    遍历所有工作表和单元格。
    返回 {"total_replacements": int, "detail": {原词: 替换次数}}
    """
    wb = load_workbook(filepath)
    detail = {}

    for old_text, new_text in rules:
        count = 0
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None and isinstance(cell.value, str):
                        if old_text in cell.value:
                            count += cell.value.count(old_text)
                            cell.value = cell.value.replace(old_text, new_text)
        if count > 0:
            detail[old_text] = count

    wb.save(filepath)
    wb.close()

    total = sum(detail.values())
    return {"total_replacements": total, "detail": detail}


# ─── XLS 替换 ─────────────────────────────────────────────

def replace_in_xls(filepath: str, rules: List[Tuple[str, str]]) -> dict:
    """
    对 xls (旧版 Excel) 文件执行批量替换。
    使用 xlrd 读取 + xlwt 写入。
    注意：格式保留有限。
    返回 {"total_replacements": int, "detail": {原词: 替换次数}}
    """
    import xlrd
    import xlwt

    rb = xlrd.open_workbook(filepath, formatting_info=True)
    wb = xlwt.Workbook()
    detail = {}

    for sheet_idx in range(rb.nsheets):
        rs = rb.sheet_by_index(sheet_idx)
        ws = wb.add_sheet(rs.name)

        for row_idx in range(rs.nrows):
            for col_idx in range(rs.ncols):
                cell_value = rs.cell_value(row_idx, col_idx)
                cell_type = rs.cell_type(row_idx, col_idx)

                if cell_type == xlrd.XL_CELL_TEXT and isinstance(cell_value, str):
                    new_value = cell_value
                    for old_text, new_text in rules:
                        cnt = new_value.count(old_text)
                        if cnt > 0:
                            detail[old_text] = detail.get(old_text, 0) + cnt
                            new_value = new_value.replace(old_text, new_text)
                    ws.write(row_idx, col_idx, new_value)
                elif cell_type == xlrd.XL_CELL_NUMBER:
                    ws.write(row_idx, col_idx, cell_value)
                elif cell_type == xlrd.XL_CELL_DATE:
                    ws.write(row_idx, col_idx, cell_value)
                elif cell_type == xlrd.XL_CELL_BOOLEAN:
                    ws.write(row_idx, col_idx, cell_value)
                else:
                    ws.write(row_idx, col_idx, cell_value)

    wb.save(filepath)

    total = sum(detail.values())
    return {"total_replacements": total, "detail": detail}


# ─── PPTX 替换 ────────────────────────────────────────────

def replace_in_pptx(filepath: str, rules: List[Tuple[str, str]]) -> dict:
    """
    对 pptx 文件执行批量替换。
    遍历所有幻灯片中的形状和文本框。
    返回 {"total_replacements": int, "detail": {原词: 替换次数}}
    """
    prs = Presentation(filepath)
    detail = {}

    def process_text_frame(text_frame):
        """处理文本框中的所有段落"""
        nonlocal detail
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                for old_text, new_text in rules:
                    if old_text in run.text:
                        count = run.text.count(old_text)
                        detail[old_text] = detail.get(old_text, 0) + count
                        run.text = run.text.replace(old_text, new_text)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                process_text_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        process_text_frame(cell.text_frame)

    prs.save(filepath)

    total = sum(detail.values())
    return {"total_replacements": total, "detail": detail}


# ─── CSV 替换 ─────────────────────────────────────────────

def replace_in_csv_file(filepath: str, rules: List[Tuple[str, str]]) -> dict:
    """
    对 csv 文件执行批量替换（当作纯文本处理）。
    返回 {"total_replacements": int, "detail": {原词: 替换次数}}
    """
    # CSV 本质上是纯文本，直接用 txt 替换即可
    return replace_in_txt(filepath, rules)


# ─── DOC 替换（旧版 Word，使用 COM 自动化）────────────────

def replace_in_doc(filepath: str, rules: List[Tuple[str, str]]) -> dict:
    """
    对 doc (旧版 Word) 文件执行批量替换。
    使用 pywin32 COM 自动化（需要安装 Word）。
    返回 {"total_replacements": int, "detail": {原词: 替换次数}}
    """
    try:
        import win32com.client
    except ImportError:
        raise RuntimeError(
            "处理 .doc 文件需要安装 pywin32 和 Microsoft Word。\n"
            "请运行: pip install pywin32"
        )

    abs_path = str(Path(filepath).resolve())
    word = None
    doc = None
    detail = {}

    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = False
        doc = word.Documents.Open(abs_path)

        for old_text, new_text in rules:
            count = 0
            # 使用 Find & Replace
            find = doc.Content.Find
            find.ClearFormatting()
            find.Replacement.ClearFormatting()
            while find.Execute(
                FindText=old_text,
                ReplaceWith=new_text,
                Replace=1,  # wdReplaceOne
                Forward=True,
                Wrap=0,      # wdFindStop
            ):
                count += 1

            if count > 0:
                detail[old_text] = count

        doc.Save()
    finally:
        if doc:
            doc.Close(False)
        if word:
            word.Quit()

    total = sum(detail.values())
    return {"total_replacements": total, "detail": detail}


# ─── PPT 替换（旧版 PowerPoint，使用 COM 自动化）──────────

def replace_in_ppt(filepath: str, rules: List[Tuple[str, str]]) -> dict:
    """
    对 ppt (旧版 PowerPoint) 文件执行批量替换。
    使用 pywin32 COM 自动化（需要安装 PowerPoint）。
    返回 {"total_replacements": int, "detail": {原词: 替换次数}}
    """
    try:
        import win32com.client
    except ImportError:
        raise RuntimeError(
            "处理 .ppt 文件需要安装 pywin32 和 Microsoft PowerPoint。\n"
            "请运行: pip install pywin32"
        )

    abs_path = str(Path(filepath).resolve())
    ppt_app = None
    presentation = None
    detail = {}

    try:
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        presentation = ppt_app.Presentations.Open(abs_path, WithWindow=False)

        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    tf = shape.TextFrame
                    if tf.HasText:
                        text = tf.TextRange.Text
                        for old_text, new_text in rules:
                            if old_text in text:
                                count = text.count(old_text)
                                detail[old_text] = detail.get(old_text, 0) + count
                                # 使用 TextRange.Replace
                                tr = tf.TextRange
                                while True:
                                    found = tr.Find(old_text)
                                    if found is None or found.Length == 0:
                                        break
                                    found.Text = new_text
                                    tr = tf.TextRange  # refresh

        presentation.Save()
    finally:
        if presentation:
            presentation.Close()
        if ppt_app:
            ppt_app.Quit()

    total = sum(detail.values())
    return {"total_replacements": total, "detail": detail}


# ─── 统一分发函数 ─────────────────────────────────────────

def replace_in_file(filepath: str, rules: List[Tuple[str, str]]) -> dict:
    """
    根据文件扩展名自动调用对应的替换函数。
    返回 {"total_replacements": int, "detail": {原词: 替换次数}}
    """
    ext = Path(filepath).suffix.lower()

    dispatch = {
        ".docx": replace_in_docx,
        ".txt":  replace_in_txt,
        ".xlsx": replace_in_xlsx,
        ".xls":  replace_in_xls,
        ".pptx": replace_in_pptx,
        ".ppt":  replace_in_ppt,
        ".csv":  replace_in_csv_file,
        ".doc":  replace_in_doc,
    }

    func = dispatch.get(ext)
    if func is None:
        raise ValueError(f"不支持的文件格式: {ext}")

    return func(filepath, rules)
